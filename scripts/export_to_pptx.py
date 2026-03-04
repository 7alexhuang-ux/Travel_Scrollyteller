import os
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

# --- Configuration ---
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# Sync with export_package.py's output directory
EXPORT_DIR_NAME = f"Golden_Journey_Presentation_{datetime.datetime.now().strftime('%Y%m%d')}"
EXPORT_DIR = os.path.join(PROJECT_ROOT, "output", EXPORT_DIR_NAME)
EXPORT_IMAGES_DIR = os.path.join(EXPORT_DIR, "images")

CLUSTERS_FILE = os.path.join(PROJECT_ROOT, "data", "location_clusters.json")
NOTES_FILE = os.path.join(PROJECT_ROOT, "notes", "緬甸_行程筆記.md")
MANUAL_LOCATIONS_FILE = os.path.join(PROJECT_ROOT, "data", "manual_cluster_locations.json")
SELECTED_PHOTOS_FILE = os.path.join(PROJECT_ROOT, "data", "selected_photos.json")
PHOTOS_SOURCE_DIR = r"C:\Project\Alex_Diary\inbox\緬甸"

OUTPUT_PPT_PATH = os.path.join(EXPORT_DIR, "presentation.pptx")

# --- Colors ---
GOLD = RGBColor(184, 134, 11)
DARK_BG = RGBColor(15, 15, 15)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(150, 150, 150)

# --- Helpers ---
def load_json(path):
    if os.path.exists(path):
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {}

def parse_notes():
    notes = {}
    if not os.path.exists(NOTES_FILE): return notes
    with open(NOTES_FILE, "r", encoding="utf-8") as f:
        content = f.read()
    pattern = r"##\s+(IMG_[A-Za-z0-9_]+\.(?:JPG|PNG|JPEG|MOV|MP4))\n(.*?)(?=\n##|\Z)"
    matches = re.finditer(pattern, content, re.DOTALL)
    for match in matches:
        filename = match.group(1).strip()
        body = match.group(2).strip()
        note_data = {"note": "", "refined_note": "", "tips": "", "subtitle": ""}
        m = re.search(r"-\s+\*\*心得\*\*:\s*(.*?)(?=\n-|\Z)", body, re.DOTALL)
        if m: note_data["note"] = m.group(1).strip()
        m = re.search(r"-\s+\*\*精修筆記\*\*:\s*(.*?)(?=\n-|\Z)", body, re.DOTALL)
        if m: note_data["refined_note"] = m.group(1).strip()
        m = re.search(r"-\s+\*\*Tips\*\*:\s*(.*?)(?=\n-|\Z)", body, re.DOTALL)
        if m: note_data["tips"] = m.group(1).strip()
        m = re.search(r"-\s+\*\*子標題\*\*:\s*(.*?)(?=\n-|\Z)", body, re.DOTALL)
        if m: note_data["subtitle"] = m.group(1).strip()
        notes[filename] = note_data
    return notes

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.text = subtitle_text
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(18)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    
    # Add tip about clicking images
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.5))
    tf3 = txBox3.text_frame
    tf3.text = "💡 提示：在簡報模式下點選照片，可直接開啟高解析度原始檔"
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(12)
    p3.font.color.rgb = GOLD
    p3.alignment = PP_ALIGN.CENTER

def add_location_slide(prs, img_filename, loc_name, note_data, is_starred):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    img_abs_path = os.path.join(EXPORT_IMAGES_DIR, img_filename)
    if not os.path.exists(img_abs_path):
        # Fallback to source if not in export yet
        img_abs_path = os.path.join(PHOTOS_SOURCE_DIR, img_filename)

    if os.path.exists(img_abs_path):
        # Image
        pic = slide.shapes.add_picture(img_abs_path, Inches(0.4), Inches(0.8), height=Inches(4.5))
        # Relative link for portability!
        pic.click_action.hyperlink.address = f"images/{img_filename}"
    
    # Text info
    tx_name = slide.shapes.add_textbox(Inches(6.6), Inches(0.5), Inches(3.2), Inches(1))
    tf_name = tx_name.text_frame
    tf_name.word_wrap = True
    tf_name.text = loc_name
    p_name = tf_name.paragraphs[0]
    p_name.font.bold = True
    p_name.font.size = Pt(22)
    p_name.font.color.rgb = GOLD
    if is_starred: p_name.text = "⭐ " + loc_name

    content = note_data.get("refined_note") or note_data.get("note") or "無備註"
    if "現場細節整理中" in content: content = note_data.get("note") or content
    
    tx_note = slide.shapes.add_textbox(Inches(6.6), Inches(1.4), Inches(3.2), Inches(3.5))
    tf_note = tx_note.text_frame
    tf_note.word_wrap = True
    tf_note.text = content
    p_note = tf_note.paragraphs[0]
    p_note.font.size = Pt(11)
    p_note.font.color.rgb = WHITE
    
    if note_data.get("tips"):
        tx_tips = slide.shapes.add_textbox(Inches(6.6), Inches(4.8), Inches(3.2), Inches(0.8))
        tf_tips = tx_tips.text_frame
        tf_tips.word_wrap = True
        tf_tips.text = "💡 " + note_data.get("tips")
        p_tips = tf_tips.paragraphs[0]
        p_tips.font.size = Pt(9)
        p_tips.font.italic = True
        p_tips.font.color.rgb = GRAY

def main():
    if not os.path.exists(EXPORT_DIR):
        print(f"Error: EXPORT_DIR {EXPORT_DIR} does not exist. Run export_package.py first.")
        return

    print(f"Generating portable PPTX into {EXPORT_DIR}...")
    clusters = load_json(CLUSTERS_FILE)
    manual_locs = load_json(MANUAL_LOCATIONS_FILE)
    selected_photos = load_json(SELECTED_PHOTOS_FILE)
    notes = parse_notes()
    starred_files = set(k for k, v in selected_photos.items() if v.get('selected', False))

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs, "緬甸建築深度考察報告", f"生成日期: {datetime.datetime.now().strftime('%Y-%m-%d')}\n包含 {len(clusters)} 個位點紀錄")

    count = 0
    for i, c in enumerate(clusters):
        sample = c['sample_file']
        loc_name = manual_locs.get(sample, f"考察位點 #{i + 1}")
        
        is_meaningful = False
        if "考察位點" not in loc_name: is_meaningful = True
        note = notes.get(sample, {})
        if note.get("note") or note.get("tips"): is_meaningful = True
        if note.get("refined_note") and "現場細節整理中" not in note.get("refined_note"): is_meaningful = True
        
        if is_meaningful:
            add_location_slide(prs, sample, loc_name, note, (sample in starred_files))
            count += 1
            if count % 50 == 0:
                print(f"Process: {count} slides...")

    prs.save(OUTPUT_PPT_PATH)
    print(f"PPTX Export complete: {OUTPUT_PPT_PATH}")
    
    # Re-zip to include the PPTX
    print("Updating ZIP package to include PPTX...")
    import shutil
    shutil.make_archive(EXPORT_DIR, 'zip', os.path.dirname(EXPORT_DIR), EXPORT_DIR_NAME)
    print(f"Final update complete: {EXPORT_DIR}.zip")

if __name__ == "__main__":
    main()
